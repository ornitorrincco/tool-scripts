import RAKE
from pptx import Presentation


def pptxGetText(filename):
    prs = Presentation(filename)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    fullText = ''

    for k in range(0, len(text_runs)):
        fullText = fullText + ' ' + text_runs[k]


filename = 'test.pptx'
fullText = pptxGetText(filename):
Rake = RAKE.Rake("SmartStopList.txt")
keywords = Rake.run(fullText)
